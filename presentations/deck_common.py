"""Shared design system for the Cyclistic conversion decks.

Design rules (v2):
  - All slides: white background.
  - CASUAL (amber) for every casual element — text markers, chips, bars, labels.
    All three casual subgroups also use CASUAL; no gradient.
  - MEMBER (blue) for every member element.
  - Body text stays INK/MUTED for readability; only markers/data elements carry group color.
  - No SURFACE card fills — cards use a hairline border or no background.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Palette ───────────────────────────────────────────────────────────────────
MEMBER   = RGBColor(0x29, 0x62, 0xA6)   # deep blue  — annual members
CASUAL   = RGBColor(0xC0, 0x6A, 0x2C)   # amber      — casual riders (all subgroups)

INK      = RGBColor(0x1C, 0x1C, 0x1C)   # primary text
MUTED    = RGBColor(0x67, 0x73, 0x82)   # secondary text
HAIRLINE = RGBColor(0xD9, 0xDF, 0xE6)   # thin separators / card borders
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

HEAD = "Trebuchet MS"
BODY = "Arial"

# 16:9 widescreen
EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def new_deck(title, author="Cyclistic Analytics"):
    prs = Presentation()
    prs.slide_width  = EMU_W
    prs.slide_height = EMU_H
    prs.core_properties.title  = title
    prs.core_properties.author = author
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _set_font(run, size, color, bold=False, font=BODY, italic=False, spacing=None):
    run.font.size        = Pt(size)
    run.font.color.rgb   = color
    run.font.bold        = bold
    run.font.italic      = italic
    run.font.name        = font
    if spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(spacing * 100)))


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, space_after=0, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (txt, opts) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap       = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment   = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, {})]
        for txt, opts in para:
            r = p.add_run()
            r.text = txt
            _set_font(r, opts.get('size', 16), opts.get('color', INK),
                      opts.get('bold', False), opts.get('font', BODY),
                      opts.get('italic', False), opts.get('spacing'))
    return tb


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.08):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width     = Pt(line_w)
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.text_frame.paragraphs[0].text = ""
    return shp


def chip(slide, x, y, color, size=0.16):
    """Small colored square — group marker next to labels."""
    return rect(slide, x, y, Inches(size), Inches(size), fill=color, rounded=True, radius=0.25)


def kicker(slide, x, y, label, color=MEMBER, num=None):
    """Section kicker: colored chip + uppercase label."""
    chip(slide, x, y + Inches(0.02), color, size=0.16)
    runs = []
    if num:
        runs.append((f"{num}   ", {'size': 12, 'color': color, 'bold': True,
                                   'font': HEAD, 'spacing': 2}))
    runs.append((label.upper(), {'size': 12, 'color': MUTED, 'bold': True,
                                 'font': HEAD, 'spacing': 2}))
    text(slide, x + Inches(0.28), y, Inches(10), Inches(0.3), [runs],
         anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, idx, total):
    text(slide, Inches(0.6), Inches(7.04), Inches(8), Inches(0.3),
         [[("Cyclistic  ·  Casual-to-Member Conversion  ·  Jun 2025 – May 2026",
            {'size': 9, 'color': MUTED, 'font': BODY})]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(slide, Inches(11.5), Inches(7.04), Inches(1.23), Inches(0.3),
         [[(f"{idx:02d} / {total:02d}", {'size': 9, 'color': MUTED, 'font': HEAD})]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def style_chart(chart, *, legend=False, legend_pos=None, value_gridlines=False,
                font_size=11):
    chart.font.name  = BODY
    chart.font.size  = Pt(font_size)
    chart.has_title  = False
    chart.has_legend = legend
    if legend and legend_pos is not None:
        from pptx.enum.chart import XL_LEGEND_POSITION
        chart.legend.position          = legend_pos
        chart.legend.include_in_layout = False
        chart.legend.font.size         = Pt(font_size)
        chart.legend.font.color.rgb    = MUTED
    try:
        va = chart.value_axis
        va.has_major_gridlines = value_gridlines
        if value_gridlines:
            gl = va.major_gridlines.format.line
            gl.color.rgb = HAIRLINE
            gl.width     = Pt(0.5)
        va.format.line.color.rgb      = HAIRLINE
        va.tick_labels.font.size      = Pt(font_size)
        va.tick_labels.font.color.rgb = MUTED
    except Exception:
        pass
    try:
        ca = chart.category_axis
        ca.has_major_gridlines        = False
        ca.format.line.color.rgb      = HAIRLINE
        ca.tick_labels.font.size      = Pt(font_size)
        ca.tick_labels.font.color.rgb = INK
    except Exception:
        pass


# ── Custom shape-based visuals ────────────────────────────────────────────────
def hsplit_bar(slide, x, y, w, h, parts, radius=0.12):
    """Proportional horizontal bar. parts = [(frac, color, label, text_color)]."""
    cx    = x
    total = sum(p[0] for p in parts)
    for k, (frac, color, label, tcol) in enumerate(parts):
        pw      = int(w * frac / total)
        rounded = (k == 0 or k == len(parts) - 1)
        rect(slide, cx, y, pw, h, fill=color, rounded=rounded, radius=radius)
        text(slide, cx, y, pw, h, [[(label, {'size': 15, 'color': tcol, 'bold': True})]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += pw


def vbars_custom(slide, x, y, w, h, items, maxval, fmt=lambda v: f"{v:.2f}",
                 barw_max=1.3, val_size=15, lab_size=13):
    """Vertical bars from rectangles. items = [(label, value, color)]."""
    n    = len(items)
    slot = int(w / n)
    barw = int(min(slot * 0.46, Inches(barw_max)))
    base = y + h
    rect(slide, x, base, w, Pt(1.5), fill=HAIRLINE)
    for i, (label, val, color) in enumerate(items):
        bh = int(h * (val / maxval)) if maxval else 0
        bx = x + slot * i + int((slot - barw) / 2)
        by = base - bh
        if bh > 0:
            rect(slide, bx, by, barw, bh, fill=color, rounded=True, radius=0.04)
        text(slide, x + slot * i, by - Inches(0.46), slot, Inches(0.4),
             [[(fmt(val), {'size': val_size, 'color': INK, 'bold': True, 'font': HEAD})]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        text(slide, x + slot * i, base + Inches(0.12), slot, Inches(0.7),
             [[(label, {'size': lab_size, 'color': INK})]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.05)
