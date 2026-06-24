"""Content + geometry QA for the decks (no renderer needed)."""
import sys
from pptx import Presentation
from pptx.util import Emu

SLIDE_W = Emu(int(13.333 * 914400))
SLIDE_H = Emu(int(7.5 * 914400))
TOL = Emu(int(0.03 * 914400))  # 0.03" tolerance


def emu_in(v):
    return round(v / 914400, 2)


def check(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    print(f"\n{'='*70}\n{path}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, {emu_in(sw)}x{emu_in(sh)} in)\n{'='*70}")
    issues = 0
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for sh_ in slide.shapes:
            if sh_.has_text_frame:
                t = " ".join(r.text for p in sh_.text_frame.paragraphs for r in p.runs).strip()
                if t:
                    texts.append(t)
            # geometry
            try:
                x, y, w, h = sh_.left, sh_.top, sh_.width, sh_.height
                if x is None:
                    continue
                if x < -TOL or y < -TOL or (x + w) > sw + TOL or (y + h) > sh + TOL:
                    issues += 1
                    nm = (texts[-1][:32] if texts else sh_.shape_type)
                    print(f"  ⚠ S{i} OUT OF BOUNDS: '{nm}' "
                          f"x={emu_in(x)} y={emu_in(y)} w={emu_in(w)} h={emu_in(h)} "
                          f"(right={emu_in(x+w)}, bottom={emu_in(y+h)})")
            except Exception:
                pass
        print(f"\n— Slide {i} —")
        for t in texts:
            print("   ", (t[:96] + "…") if len(t) > 96 else t)
    print(f"\n>>> geometry issues: {issues}")
    return issues


if __name__ == "__main__":
    total = 0
    for p in sys.argv[1:]:
        total += check(p)
    sys.exit(0)
