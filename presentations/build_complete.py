"""Complete deck (18 slides) — white background, strict group color assignment."""
from deck_common import *
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

prs   = new_deck("Cyclistic — Casual-to-Member Conversion (Full Analysis)")
TOTAL = 18
LX, RW = Inches(0.6), Inches(12.13)


def title_block(s, kick, num, color, title, sub=None, tsize=29):
    kicker(s, LX, Inches(0.55), kick, color, num)
    text(s, LX, Inches(0.92), RW, Inches(0.95),
         [[(title, {'size': tsize, 'color': INK, 'bold': True, 'font': HEAD})]],
         anchor=MSO_ANCHOR.TOP)
    if sub:
        text(s, LX, Inches(1.74), Inches(11.4), Inches(0.5),
             [[(sub, {'size': 14, 'color': MUTED})]])


def stat(s, x, y, big, label, color, bw=2.4):
    text(s, Inches(x), Inches(y), Inches(bw), Inches(0.7),
         [[(big, {'size': 30, 'color': color, 'bold': True, 'font': HEAD})]])
    text(s, Inches(x), Inches(y + 0.62), Inches(bw), Inches(0.45),
         [[(label, {'size': 12, 'color': MUTED})]])


# ── 1 · Title ────────────────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
rect(s, LX, Inches(2.35), Inches(0.55), Inches(0.07), fill=CASUAL)
text(s, LX, Inches(2.55), Inches(12.5), Inches(2.1), [
    [("Casual-to-member conversion:", {'size': 38, 'color': INK, 'bold': True, 'font': HEAD})],
    [("a behavioural segmentation", {'size': 38, 'color': CASUAL, 'bold': True, 'font': HEAD})],
], line_spacing=1.07)
text(s, LX, Inches(4.7), Inches(11.5), Inches(0.5),
     [[("Full analysis: behaviour, method, segments, and a conversion plan",
        {'size': 16, 'color': MUTED})]])
text(s, LX, Inches(5.5), Inches(11.5), Inches(0.35),
     [[("Chicago bike-share  ·  5.84 M rides  ·  12 months (Jun 2025 – May 2026)  ·  Divvy public data",
        {'size': 11, 'color': MUTED, 'font': HEAD, 'spacing': 1})]])
chip(s, LX, Inches(6.38), CASUAL)
text(s, Inches(0.92), Inches(6.31), Inches(2.2), Inches(0.35),
     [[("Casual riders", {'size': 12, 'color': CASUAL})]], anchor=MSO_ANCHOR.MIDDLE)
chip(s, Inches(2.6), Inches(6.38), MEMBER)
text(s, Inches(2.92), Inches(6.31), Inches(2.8), Inches(0.35),
     [[("Annual members", {'size': 12, 'color': MEMBER})]], anchor=MSO_ANCHOR.MIDDLE)

# ── 2 · Contents ─────────────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "Contents", None, MEMBER, "What this deck covers")
items = [
    ("01", "The question",         "Why conversion, and what we ask of the data"),
    ("02", "How casual differs",   "Duration, timing and the weekday rhythm"),
    ("03", "Method",               "Features, clustering, and choosing K"),
    ("04", "The three segments",   "Profiles ranked by member similarity"),
    ("05", "Priority & actions",   "Where to focus and concrete recommendations"),
    ("06", "Limits & next steps",  "What the analysis can and cannot say"),
]
colx = [0.6, 6.8]
for k, (num, h, d) in enumerate(items):
    col = k // 3
    row = k % 3
    x  = colx[col]
    yy = 2.15 + row * 1.45
    text(s, Inches(x), Inches(yy), Inches(1.1), Inches(0.7),
         [[(num, {'size': 24, 'color': MEMBER, 'bold': True, 'font': HEAD})]])
    text(s, Inches(x + 1.0), Inches(yy), Inches(5.0), Inches(0.45),
         [[(h, {'size': 17, 'color': INK, 'bold': True})]])
    text(s, Inches(x + 1.0), Inches(yy + 0.42), Inches(5.0), Inches(0.7),
         [[(d, {'size': 12.5, 'color': MUTED})]], line_spacing=1.1)
footer(s, 2, TOTAL)

# ── 3 · Business question ─────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "The question", "01", CASUAL, "Which casual riders are worth converting?")
text(s, LX, Inches(1.85), Inches(11.7), Inches(0.85),
     [[("Cyclistic has two user types — ", {'size': 15, 'color': INK}),
       ("casual riders", {'size': 15, 'color': CASUAL, 'bold': True}),
       (" (single / day pass) and ", {'size': 15, 'color': INK}),
       ("annual members", {'size': 15, 'color': MEMBER, 'bold': True}),
       (". Converting the right casual riders is the cheapest growth lever.",
        {'size': 15, 'color': INK})]], line_spacing=1.2)
qs = [("01", "Do casual riders behave differently from members?"),
      ("02", "Are casual riders one group — or several distinct profiles?"),
      ("03", "Which profiles look most like members, and are big enough to target?")]
y = 3.0
for num, q in qs:
    rect(s, LX, Inches(y), Inches(0.07), Inches(0.78), fill=CASUAL)
    text(s, Inches(0.9), Inches(y), Inches(1.1), Inches(0.78),
         [[(num, {'size': 22, 'color': CASUAL, 'bold': True, 'font': HEAD})]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.95), Inches(y), Inches(10.2), Inches(0.78),
         [[(q, {'size': 15.5, 'color': INK})]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(1.95), Inches(y + 0.8), Inches(10.2), Pt(0.5), fill=HAIRLINE)
    y += 0.98
# Note box (hairline border)
rect(s, LX, Inches(6.08), RW, Inches(0.62),
     fill=None, line=HAIRLINE, line_w=0.75, rounded=True, radius=0.1)
text(s, Inches(0.9), Inches(6.08), Inches(11.6), Inches(0.62),
     [[("Note  ", {'size': 12, 'color': CASUAL, 'bold': True}),
       ("Segmentation is at ride level — no persistent user ID — so clusters describe patterns "
        "of behaviour, not individual people.", {'size': 12, 'color': INK})]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3, TOTAL)

# ── 4 · Data & scope ──────────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "The data", "02", CASUAL, "12 months of Chicago bike-share trips")
stat(s, 0.6,  2.15, "5.84 M", "rides analysed", INK)
stat(s, 3.4,  2.15, "12",     "months (rolling window)", INK)
stat(s, 5.7,  2.15, "13",     "fields per ride", INK)
stat(s, 8.0,  2.15, "Divvy",  "public open data", MEMBER)
rect(s, LX, Inches(3.45), RW, Pt(1), fill=HAIRLINE)
text(s, LX, Inches(3.7), Inches(11.8), Inches(0.4),
     [[("The split: of every 100 rides, ", {'size': 15, 'color': INK}),
       ("~36 are casual", {'size': 15, 'color': CASUAL, 'bold': True}),
       (".", {'size': 15, 'color': INK})]])
hsplit_bar(s, LX, Inches(4.25), RW, Inches(0.72),
           [(64.1, MEMBER, "Members  64.1%", WHITE),
            (35.9, CASUAL, "Casual  35.9%",  WHITE)])
facts = ["Cleaned in a separate, documented pipeline (notebook 01).",
         "Rides with missing end-GPS and timestamp errors removed (<0.2%).",
         "Free-floating e-bikes legitimately have no end-station — kept by design.",
         "Features engineered: duration, distance, speed, weekend, commute window, e-bike, season."]
y = 5.35
for f in facts:
    chip(s, LX, Inches(y + 0.04), MEMBER, size=0.14)
    text(s, Inches(0.92), Inches(y), Inches(11.6), Inches(0.35),
         [[(f, {'size': 13, 'color': INK})]])
    y += 0.33
footer(s, 4, TOTAL)

# ── 5 · How casual differs ────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "First look", "02", CASUAL, "Casual and member riders use the service differently")
cd = CategoryChartData()
cd.categories = ["Median ride (min)", "Rides in commute window (%)"]
cd.add_series("Casual",  (11, 17))
cd.add_series("Members", (9,  28))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(2.2),
                        Inches(7.2), Inches(4.3), cd)
ch = gf.chart
ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb = CASUAL
ch.series[1].format.fill.solid(); ch.series[1].format.fill.fore_color.rgb = MEMBER
style_chart(ch, legend=True, legend_pos=XL_LEGEND_POSITION.TOP)
ch.plots[0].has_data_labels = True
dl = ch.plots[0].data_labels
dl.number_format = '0'; dl.number_format_is_linked = False
dl.font.size = Pt(12); dl.font.bold = True; dl.font.color.rgb = INK
dl.position = 3
ch.value_axis.visible           = False
ch.value_axis.has_major_gridlines = False
findings = [("Longer casual rides",
             "Casual median is 11 min vs 9 for members — and the casual mean is pulled far "
             "higher by long recreational trips."),
            ("Members commute more",
             "28% of member rides fall in commute windows vs 17% for casual."),
            ("Casual skews to weekends",
             "Recreation, not routine, dominates casual demand.")]
y = 2.35
for h, d in findings:
    chip(s, Inches(8.1), Inches(y + 0.05), CASUAL, size=0.16)
    text(s, Inches(8.42), Inches(y - 0.04), Inches(4.3), Inches(0.4),
         [[(h, {'size': 15, 'color': INK, 'bold': True})]])
    text(s, Inches(8.42), Inches(y + 0.33), Inches(4.3), Inches(0.95),
         [[(d, {'size': 12, 'color': MUTED})]], line_spacing=1.1)
    y += 1.4
footer(s, 5, TOTAL)

# ── 6 · Weekday rhythm ────────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "The weekday rhythm", "02", CASUAL, "The morning peak is where members stand alone")
hours  = [str(h) for h in range(6, 23)]
member = [12, 30, 78, 100, 60, 45, 42, 48, 55, 62, 78, 100, 92, 60, 38, 26, 18]
casual = [6,  12, 26, 34,  30, 34, 40, 48, 56, 64, 78, 92,  88, 74, 58, 44, 32]
cd = CategoryChartData()
cd.categories = hours
cd.add_series("Casual",  casual)
cd.add_series("Members", member)
gf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.55), Inches(2.25),
                        Inches(8.0), Inches(4.1), cd)
ch = gf.chart
ch.series[0].format.line.color.rgb = CASUAL; ch.series[0].format.line.width = Pt(2.75)
ch.series[1].format.line.color.rgb = MEMBER; ch.series[1].format.line.width = Pt(2.75)
ch.series[0].smooth = True; ch.series[1].smooth = True
style_chart(ch, legend=True, legend_pos=XL_LEGEND_POSITION.TOP)
ch.value_axis.visible             = False
ch.value_axis.has_major_gridlines = False
# Annotation panel: hairline bordered
rect(s, Inches(8.8), Inches(2.55), Inches(3.95), Inches(3.5),
     fill=None, line=HAIRLINE, line_w=0.75, rounded=True, radius=0.05)
text(s, Inches(9.1), Inches(2.8), Inches(3.4), Inches(0.38),
     [[("WHAT TO SEE", {'size': 11, 'color': CASUAL, 'bold': True, 'font': HEAD, 'spacing': 2})]])
text(s, Inches(9.1), Inches(3.22), Inches(3.4), Inches(2.65),
     [[("Members show two sharp peaks — morning and evening commute.", {'size': 13, 'color': INK})],
      [("Casual riders ", {'size': 13, 'color': INK}),
       ("miss the morning peak", {'size': 13, 'color': CASUAL, 'bold': True}),
       (" almost entirely, but share the evening one.", {'size': 13, 'color': INK})],
      [("The AM weekday peak is the cleanest signal of member-like behaviour.",
        {'size': 13, 'color': INK})]],
     line_spacing=1.18, space_after=8)
text(s, LX, Inches(6.5), RW, Inches(0.3),
     [[("Hour of departure (weekdays). Schematic of the documented pattern — relative, not to scale.",
        {'size': 10, 'color': MUTED, 'italic': True})]])
footer(s, 6, TOTAL)

# ── 7 · Method ────────────────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "Method", "03", MEMBER, "How the segments were found")
steps = [
    ("Engineer features",  "8 behavioural features per ride: duration, distance, speed, weekend, commute window, time-of-day, season, e-bike."),
    ("Cluster casual rides","MiniBatch K-Means groups casual rides into behavioural profiles — letting the data define the segments."),
    ("Benchmark vs members","Members are clustered the same way; each casual segment is scored by distance to the member profile."),
    ("Prioritise",          "Combined score = similarity × ride share — balancing behavioural fit against audience size."),
]
x  = 0.6
cw = 2.96
for i, (h, d) in enumerate(steps):
    rect(s, Inches(x), Inches(2.3), Inches(cw), Inches(3.7),
         fill=None, line=HAIRLINE, line_w=0.75, rounded=True, radius=0.05)
    text(s, Inches(x + 0.28), Inches(2.55), Inches(1.0), Inches(0.7),
         [[(f"{i+1}", {'size': 30, 'color': MEMBER, 'bold': True, 'font': HEAD})]])
    text(s, Inches(x + 0.28), Inches(3.4), Inches(cw - 0.5), Inches(0.7),
         [[(h, {'size': 15, 'color': INK, 'bold': True})]], line_spacing=1.0)
    text(s, Inches(x + 0.28), Inches(4.15), Inches(cw - 0.5), Inches(1.7),
         [[(d, {'size': 12, 'color': MUTED})]], line_spacing=1.15)
    x += cw + 0.12
footer(s, 7, TOTAL)

# ── 8 · Choosing K ────────────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "Validation", "03", MEMBER, "Why three segments")
text(s, LX, Inches(1.9), Inches(11.7), Inches(1.0),
     [[("K was chosen empirically — inertia and silhouette compared across K = 2, 3, 4. "
        "The inertia curve bends at K=3 (the 2→3 drop is more than twice the 3→4 drop). "
        "K=2 would only split weekday vs. weekend — too coarse to act on.",
        {'size': 15, 'color': INK})]], line_spacing=1.25)
boxes = [("K = 2", "Too coarse",   "Just weekday vs. weekend — not actionable.", False),
         ("K = 3", "Chosen",        "Elbow supports it; three profiles a team can act on.", True),
         ("K = 4", "Marginal gain", "Little extra separation for added complexity.", False)]
x = 0.6
for kval, label, body, chosen in boxes:
    lc = CASUAL if chosen else HAIRLINE
    rect(s, Inches(x), Inches(3.25), Inches(3.95), Inches(2.4),
         fill=None, line=lc, line_w=(2.0 if chosen else 0.75), rounded=True, radius=0.05)
    text(s, Inches(x + 0.3), Inches(3.5), Inches(3.4), Inches(0.6),
         [[(kval, {'size': 24, 'color': (CASUAL if chosen else INK), 'bold': True, 'font': HEAD})]])
    text(s, Inches(x + 0.3), Inches(4.18), Inches(3.4), Inches(0.4),
         [[(label.upper(), {'size': 12, 'color': CASUAL if chosen else MUTED,
                            'bold': True, 'font': HEAD, 'spacing': 1})]])
    text(s, Inches(x + 0.3), Inches(4.6), Inches(3.4), Inches(0.95),
         [[(body, {'size': 12.5, 'color': INK if chosen else MUTED})]], line_spacing=1.15)
    x += 3.95 + 0.14
# Honesty note
rect(s, LX, Inches(5.85), RW, Inches(0.72),
     fill=None, line=HAIRLINE, line_w=0.75, rounded=True, radius=0.08)
text(s, Inches(0.9), Inches(5.85), Inches(11.6), Inches(0.72),
     [[("Note  ", {'size': 12, 'color': CASUAL, 'bold': True}),
       ("Silhouette scores are modest in absolute terms — segments are soft behavioural regions, "
        "not sharply separated groups. Useful for targeting, not hard borders.",
        {'size': 12, 'color': INK})]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
footer(s, 8, TOTAL)

# ── 9 · Three segments overview ──────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "The segments", "04", CASUAL, "Three profiles, ranked by member similarity")
cards = [
    ("E-Bike Riders",         "52.1%", "0.65", "HIGH",
     "Short, purposeful trips on weekdays. Already ride like members."),
    ("Commute E-Bike Riders", "17.2%", "0.49", "MEDIUM",
     "100% rush-hour rides; ~1 in 4 ends outside a dock."),
    ("Classic Leisure Riders","30.8%", "0.00", "LOW",
     "Longer recreational rides, weekend-skewed."),
]
cw, gap, x = 3.95, 0.14, 0.6
for name, share, sim, pri, desc in cards:
    rect(s, Inches(x), Inches(2.2), Inches(cw), Inches(4.3),
         fill=None, line=HAIRLINE, line_w=0.75, rounded=True, radius=0.05)
    rect(s, Inches(x), Inches(2.2), Inches(cw), Inches(0.1), fill=CASUAL)
    text(s, Inches(x + 0.3), Inches(2.5), Inches(cw - 0.5), Inches(0.75),
         [[(name, {'size': 17, 'color': INK, 'bold': True, 'font': HEAD})]])
    rect(s, Inches(x + 0.3), Inches(3.3), Inches(1.65), Inches(0.4),
         fill=CASUAL, rounded=True, radius=0.5)
    text(s, Inches(x + 0.3), Inches(3.3), Inches(1.65), Inches(0.4),
         [[(pri + " PRIORITY", {'size': 10, 'color': WHITE, 'bold': True,
                                'font': HEAD, 'spacing': 1})]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(x + 0.3), Inches(4.0), Inches(1.7), Inches(0.95),
         [[(share, {'size': 25, 'color': CASUAL, 'bold': True, 'font': HEAD})],
          [("of casual rides", {'size': 11, 'color': MUTED})]])
    text(s, Inches(x + 2.2), Inches(4.0), Inches(1.6), Inches(0.95),
         [[(sim, {'size': 25, 'color': MEMBER, 'bold': True, 'font': HEAD})],
          [("similarity", {'size': 11, 'color': MUTED})]])
    text(s, Inches(x + 0.3), Inches(5.1), Inches(cw - 0.55), Inches(1.2),
         [[(desc, {'size': 13, 'color': INK})]], line_spacing=1.15)
    x += cw + gap
footer(s, 9, TOTAL)


# ── Segment deep-dive template (slides 10-12) ─────────────────────────────────
def segment_slide(idx, name, pri, share, sim, lead, traits, action):
    s = blank(prs); bg(s, WHITE)
    title_block(s, f"Segment · {pri}", "04", CASUAL, name)
    # Left stats panel: hairline border
    rect(s, LX, Inches(2.2), Inches(3.7), Inches(4.3),
         fill=None, line=HAIRLINE, line_w=0.75, rounded=True, radius=0.05)
    rect(s, LX, Inches(2.2), Inches(3.7), Inches(0.1), fill=CASUAL)
    rect(s, Inches(0.9), Inches(2.55), Inches(1.8), Inches(0.42),
         fill=CASUAL, rounded=True, radius=0.5)
    text(s, Inches(0.9), Inches(2.55), Inches(1.8), Inches(0.42),
         [[(pri + " PRIORITY", {'size': 10, 'color': WHITE, 'bold': True,
                                'font': HEAD, 'spacing': 1})]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.9), Inches(3.3), Inches(3.1), Inches(1.0),
         [[(share, {'size': 40, 'color': CASUAL, 'bold': True, 'font': HEAD})],
          [("of all casual rides", {'size': 12, 'color': MUTED})]])
    text(s, Inches(0.9), Inches(4.75), Inches(3.1), Inches(1.0),
         [[(sim, {'size': 40, 'color': MEMBER, 'bold': True, 'font': HEAD})],
          [("member similarity (0–1)", {'size': 12, 'color': MUTED})]])
    # Right: lead statement
    text(s, Inches(4.6), Inches(2.2), Inches(8.1), Inches(0.95),
         [[(lead, {'size': 17, 'color': INK, 'bold': True})]], line_spacing=1.2)
    # Traits with amber bullets
    y = 3.35
    for t in traits:
        chip(s, Inches(4.6), Inches(y + 0.04), CASUAL, size=0.14)
        text(s, Inches(4.88), Inches(y), Inches(7.8), Inches(0.5),
             [[(t, {'size': 14, 'color': INK})]], line_spacing=1.1)
        y += 0.52
    # Action bar: amber left accent
    rect(s, Inches(4.6), Inches(5.55), Inches(0.08), Inches(0.95), fill=CASUAL)
    text(s, Inches(4.83), Inches(5.62), Inches(1.4), Inches(0.35),
         [[("ACTION", {'size': 11, 'color': CASUAL, 'bold': True, 'font': HEAD, 'spacing': 2})]])
    text(s, Inches(4.83), Inches(5.97), Inches(7.8), Inches(0.5),
         [[(action, {'size': 13.5, 'color': INK})]], line_spacing=1.15)
    footer(s, idx, TOTAL)


# ── 10 · E-Bike Riders ───────────────────────────────────────────────────────
segment_slide(10, "E-Bike Riders", "HIGH", "52.1%", "0.65",
    "The largest casual segment — and the one that already rides like a member.",
    ["Short, purposeful trips concentrated on weekdays.",
     "Predominantly electric bikes, used for getting somewhere — not sightseeing.",
     "Closest of all casual segments to the member behavioural profile.",
     "Size and fit align: the obvious lead target."],
    "Lead the annual-membership campaign here. Trigger on short weekday e-bike rides.")

# ── 11 · Commute E-Bike Riders ───────────────────────────────────────────────
segment_slide(11, "Commute E-Bike Riders", "MEDIUM", "17.2%", "0.49",
    "A rush-hour segment with a physical-infrastructure story attached.",
    ["100% of rides fall in commute windows; ~70% are on e-bikes.",
     "About 1 in 4 rides ends outside a docking station (free-floating).",
     "A free-floating park-up means no guaranteed bike near home next morning.",
     "Smaller and a step further from members than E-Bike Riders."],
    "Pair a membership offer with a docking pilot in residential areas where rides end.")

# ── 12 · Classic Leisure Riders ──────────────────────────────────────────────
segment_slide(12, "Classic Leisure Riders", "LOW", "30.8%", "0.00",
    "Recreational riders whose behaviour is furthest from membership.",
    ["Longer rides, classic bikes, concentrated on weekends.",
     "Use the service for leisure, not routine transport.",
     "Most distant from the member profile in this dataset (similarity 0.00).",
     "A sizeable group — but a poor fit for an annual-membership pitch."],
    "Do not push annual membership. Offer a weekend or seasonal pass instead.")

# ── 13 · Prioritisation ──────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "Prioritisation", "05", CASUAL, "Balancing fit against audience size")
text(s, LX, Inches(1.85), Inches(11.8), Inches(0.55),
     [[("Priority = member similarity × share of casual rides. "
        "Both dimensions point at the same segment.",
        {'size': 14.5, 'color': INK})]], line_spacing=1.2)
vbars_custom(s, Inches(0.7), Inches(2.95), Inches(7.2), Inches(3.05),
             [("E-Bike\nRiders",         0.339, CASUAL),
              ("Commute\nE-Bike Riders",  0.084, CASUAL),
              ("Classic\nLeisure Riders", 0.0,   CASUAL)],
             maxval=0.40, fmt=lambda v: f"{v:.2f}", val_size=16)
# Callout: amber left bar
rect(s, Inches(8.35), Inches(2.65), Inches(0.08), Inches(3.85), fill=CASUAL)
text(s, Inches(8.58), Inches(2.85), Inches(4.5), Inches(0.38),
     [[("THE TAKEAWAY", {'size': 11, 'color': CASUAL, 'bold': True, 'font': HEAD, 'spacing': 2})]])
text(s, Inches(8.58), Inches(3.3), Inches(4.5), Inches(3.1),
     [[("E-Bike Riders score ~4× the next segment.", {'size': 15, 'color': INK, 'bold': True})],
      [("They are both the largest casual segment and the most member-like — "
        "so volume and fit reinforce each other rather than trade off.",
        {'size': 13.5, 'color': MUTED})],
      [("Classic Leisure Riders score ~0: large, but behaviourally too far for membership.",
        {'size': 13.5, 'color': MUTED})]],
     line_spacing=1.22, space_after=10)
footer(s, 13, TOTAL)

# ── 14 · Recommendations ──────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "Recommendations", "05", CASUAL, "Four moves, matched to behaviour")
recs = [
    (CASUAL, "01", "Lead with E-Bike Riders",
     "Push annual membership to this segment first; message commute-cost savings vs. pay-per-ride."),
    (CASUAL, "02", "Sell leisure riders a different product",
     "Offer Classic Leisure Riders a weekend or seasonal pass — not annual membership."),
    (CASUAL, "03", "Test the infrastructure lever",
     "Pilot added docking where PM commute e-bike rides end; it may remove a real conversion barrier."),
    (MUTED,  "04", "Prove it before scaling",
     "Run a controlled A/B campaign and measure real conversion before committing budget."),
]
y = 2.15
for c, num, head, body in recs:
    rect(s, LX, Inches(y), Inches(0.08), Inches(0.98), fill=c)
    text(s, Inches(0.9), Inches(y + 0.12), Inches(0.95), Inches(0.76),
         [[(num, {'size': 22, 'color': c, 'bold': True, 'font': HEAD})]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.95), Inches(y + 0.12), Inches(10.7), Inches(0.38),
         [[(head, {'size': 16, 'color': INK, 'bold': True})]])
    text(s, Inches(1.95), Inches(y + 0.5), Inches(10.7), Inches(0.45),
         [[(body, {'size': 12.5, 'color': MUTED})]], line_spacing=1.1)
    rect(s, Inches(1.95), Inches(y + 1.0), Inches(10.7), Pt(0.5), fill=HAIRLINE)
    y += 1.12
footer(s, 14, TOTAL)

# ── 15 · Action map ───────────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "Action map", "05", CASUAL, "The right offer for each segment")
colx = [0.6,  3.7,  6.55, 9.95]
colw = [3.0,  2.75, 3.3,  2.78]
ytab = 2.3
for j, h in enumerate(["Segment", "Offer", "Channel signal", "Expectation"]):
    text(s, Inches(colx[j]), Inches(ytab), Inches(colw[j]), Inches(0.35),
         [[(h.upper(), {'size': 11, 'color': MUTED, 'bold': True,
                        'font': HEAD, 'spacing': 1})]])
rect(s, LX, Inches(ytab + 0.42), RW, Pt(1.2), fill=HAIRLINE)
rowdata = [
    ("E-Bike Riders",         "Annual membership",         "Short weekday e-bike trips",     "Lead target — scale after test"),
    ("Commute E-Bike Riders", "Membership + docking pilot", "Rush-hour, free-floating ends",  "Watch — infra may unlock"),
    ("Classic Leisure Riders","Weekend / seasonal pass",    "Long weekend rides",             "Don't push annual"),
]
y = ytab + 0.62
for cells in rowdata:
    chip(s, Inches(colx[0] + 0.18), Inches(y + 0.47), CASUAL, size=0.18)
    text(s, Inches(colx[0] + 0.5), Inches(y), Inches(colw[0] - 0.5), Inches(1.12),
         [[(cells[0], {'size': 14, 'color': INK, 'bold': True})]], anchor=MSO_ANCHOR.MIDDLE)
    for j in range(1, 4):
        text(s, Inches(colx[j]), Inches(y), Inches(colw[j] - 0.15), Inches(1.12),
             [[(cells[j], {'size': 12.5, 'color': INK if j == 1 else MUTED})]],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    rect(s, LX, Inches(y + 1.14), RW, Pt(0.5), fill=HAIRLINE)
    y += 1.28
footer(s, 15, TOTAL)

# ── 16 · Roadmap ──────────────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "Roadmap", "05", CASUAL, "Test before you scale")
phases = [
    ("Phase 1", "Pilot",  CASUAL,
     "Controlled A/B membership offer to E-Bike Riders. Hold-out group; measure real conversion, not similarity."),
    ("Phase 2", "Learn",  MEMBER,
     "Validate the infrastructure signal — map free-floating endpoints; pilot docking in dock-sparse residential areas."),
    ("Phase 3", "Scale",  CASUAL,
     "Roll out what converts. Launch a separate weekend/seasonal product for Classic Leisure Riders."),
]
x  = 0.6
cw = 3.95
for i, (ph, h, c, d) in enumerate(phases):
    rect(s, Inches(x), Inches(2.4), Inches(cw), Inches(3.6),
         fill=None, line=HAIRLINE, line_w=0.75, rounded=True, radius=0.05)
    text(s, Inches(x + 0.3), Inches(2.65), Inches(cw - 0.6), Inches(0.38),
         [[(ph.upper(), {'size': 11, 'color': c, 'bold': True, 'font': HEAD, 'spacing': 2})]])
    text(s, Inches(x + 0.3), Inches(3.05), Inches(cw - 0.6), Inches(0.6),
         [[(h, {'size': 22, 'color': INK, 'bold': True, 'font': HEAD})]])
    text(s, Inches(x + 0.3), Inches(3.8), Inches(cw - 0.55), Inches(2.0),
         [[(d, {'size': 13, 'color': MUTED})]], line_spacing=1.2)
    if i < 2:
        text(s, Inches(x + cw - 0.02), Inches(3.9), Inches(0.4), Inches(0.5),
             [[("→", {'size': 22, 'color': HAIRLINE, 'bold': True})]])
    x += cw + 0.14
footer(s, 16, TOTAL)

# ── 17 · Limitations ──────────────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
title_block(s, "Read with care", "06", CASUAL, "Limitations of the analysis")
lims = [
    ("No user ID",              "Segmentation is ride-level; one person can appear in several segments. Sizes count rides, not people."),
    ("Similarity ≠ conversion", "It is a structural indicator. Price, alternatives and habit — not in this data — drive real decisions."),
    ("Soft clusters",           "Modest silhouette scores: segments are behavioural regions, not sharply separated groups."),
    ("Commute is a proxy",      "Departure time approximates intent; commute windows also carry a high mechanical base rate (~28%)."),
    ("~33% missing station data","GPS used instead; some station-based extensions aren't available."),
    ("Straight-line distance",  "Haversine underestimates true cycling distance, so speed is approximate."),
]
for i, (h, d) in enumerate(lims):
    col = i % 2
    row = i // 2
    xx  = 0.6 + col * 6.25
    yy  = 2.15 + row * 1.5
    chip(s, Inches(xx), Inches(yy + 0.06), CASUAL, size=0.16)
    text(s, Inches(xx + 0.32), Inches(yy - 0.03), Inches(5.7), Inches(0.4),
         [[(h, {'size': 15, 'color': INK, 'bold': True})]])
    text(s, Inches(xx + 0.32), Inches(yy + 0.34), Inches(5.7), Inches(1.0),
         [[(d, {'size': 12.5, 'color': MUTED})]], line_spacing=1.12)
footer(s, 17, TOTAL)

# ── 18 · Next steps + closing ─────────────────────────────────────────────────
s = blank(prs); bg(s, WHITE)
kicker(s, LX, Inches(0.6), "To sharpen the next round", MEMBER)
text(s, LX, Inches(1.0), Inches(12), Inches(0.75),
     [[("What would make this sharper", {'size': 28, 'color': INK, 'bold': True, 'font': HEAD})]])
nexts = [("User IDs",          "Segment by person, not by ride"),
         ("Survey data",       "Validate intent behind the behaviour"),
         ("Conversion history","Calibrate similarity against real outcomes")]
x = 0.6
for h, d in nexts:
    rect(s, Inches(x), Inches(2.1), Inches(3.95), Inches(1.7),
         fill=None, line=HAIRLINE, line_w=0.75, rounded=True, radius=0.06)
    chip(s, Inches(x + 0.3), Inches(2.45), MEMBER, size=0.18)
    text(s, Inches(x + 0.3), Inches(2.75), Inches(3.4), Inches(0.5),
         [[(h, {'size': 18, 'color': INK, 'bold': True, 'font': HEAD})]])
    text(s, Inches(x + 0.3), Inches(3.25), Inches(3.4), Inches(0.5),
         [[(d, {'size': 13, 'color': MUTED})]], line_spacing=1.1)
    x += 3.95 + 0.14
rect(s, LX, Inches(4.35), Inches(0.55), Inches(0.07), fill=CASUAL)
text(s, LX, Inches(4.55), Inches(12), Inches(0.5),
     [[("THE DECISION", {'size': 14, 'color': CASUAL, 'bold': True, 'font': HEAD, 'spacing': 2})]])
text(s, LX, Inches(5.1), Inches(12), Inches(1.55), [
    [("Run a controlled membership campaign on ",
      {'size': 26, 'color': INK, 'bold': True, 'font': HEAD}),
     ("E-Bike Riders", {'size': 26, 'color': CASUAL, 'bold': True, 'font': HEAD})],
    [("— then scale on results.", {'size': 26, 'color': INK, 'bold': True, 'font': HEAD})],
], line_spacing=1.1)
text(s, Inches(8.0), Inches(6.82), Inches(5.1), Inches(0.35),
     [[("Cyclistic Analytics  ·  Divvy public data  ·  Jun 2025 – May 2026",
        {'size': 10, 'color': MUTED, 'font': HEAD, 'spacing': 1})]],
     align=PP_ALIGN.RIGHT)

out = os.path.join(os.path.dirname(__file__), "Cyclistic_Conversion_Full.pptx")
prs.save(out)
print("Saved:", out, "·", TOTAL, "slides")
