"""Render every slide to PNG using macOS QuickLook (no LibreOffice needed).

QuickLook only thumbnails the FIRST slide of a .pptx, so for each slide i we
write a copy of the deck with slide i moved to the front, then qlmanage it.
"""
import sys, os, copy, subprocess, shutil
from pptx import Presentation


def render(path, outdir):
    base = os.path.splitext(os.path.basename(path))[0]
    tmpdir = os.path.join(outdir, "_tmp_" + base)
    shutil.rmtree(tmpdir, ignore_errors=True)
    os.makedirs(tmpdir, exist_ok=True)
    os.makedirs(outdir, exist_ok=True)

    prs = Presentation(path)
    n = len(prs.slides._sldIdLst)
    pngs = []
    for i in range(n):
        prs2 = Presentation(path)
        lst = prs2.slides._sldIdLst
        ids = list(lst)
        # move slide i to front
        target = ids[i]
        lst.remove(target)
        lst.insert(0, target)
        tmp_pptx = os.path.join(tmpdir, f"{base}_{i+1:02d}.pptx")
        prs2.save(tmp_pptx)
        subprocess.run(["qlmanage", "-t", "-s", "1600", "-o", tmpdir, tmp_pptx],
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        src_png = tmp_pptx + ".png"
        dst_png = os.path.join(outdir, f"{base}-{i+1:02d}.png")
        if os.path.exists(src_png):
            shutil.move(src_png, dst_png)
            pngs.append(dst_png)
        else:
            print(f"  ! slide {i+1}: no thumbnail")
    shutil.rmtree(tmpdir, ignore_errors=True)
    print(f"{base}: {len(pngs)}/{n} slides rendered → {outdir}")
    for p in pngs:
        print("  ", p)


if __name__ == "__main__":
    outdir = sys.argv[1]
    for path in sys.argv[2:]:
        render(path, outdir)
